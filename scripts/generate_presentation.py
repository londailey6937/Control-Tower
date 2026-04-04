#!/usr/bin/env python3
"""Generate 510k Bridge sales presentations (EN + CN) with screenshots."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCREENSHOTS = os.path.expanduser("~/arch-medical/public/screenshots")
OUT_DIR = os.path.expanduser("~/arch-medical/downloads")

# Brand colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x1A, 0x24, 0x3B)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
GREEN = RGBColor(0x34, 0xD3, 0x99)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title_text(slide, text, font_size=40, top=Inches(0.4), font_name="Calibri"):
    return add_text_box(slide, Inches(0.8), top, Inches(11.7), Inches(0.8),
                        text, font_size, ACCENT, True, PP_ALIGN.LEFT, font_name)


def add_subtitle_text(slide, text, font_size=20, top=Inches(1.2), font_name="Calibri"):
    return add_text_box(slide, Inches(0.8), top, Inches(11.7), Inches(0.6),
                        text, font_size, LIGHT_GRAY, False, PP_ALIGN.LEFT, font_name)


def add_screenshot(slide, filename, left, top, width, suffix=""):
    """Add screenshot; uses -cn suffix for CN version if available."""
    base, ext = os.path.splitext(filename)
    path = os.path.join(SCREENSHOTS, f"{base}{suffix}{ext}")
    if not os.path.exists(path):
        path = os.path.join(SCREENSHOTS, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width)


def add_card(slide, left, top, width, height, icon, title, body,
             font_name="Calibri"):
    """Add a rounded-rect card with icon, title, and body text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(12)
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)

    p0 = tf.paragraphs[0]
    p0.text = icon
    p0.font.size = Pt(28)
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = font_name
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = LIGHT_GRAY
    p2.font.name = font_name
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)


def add_screenshot_card(slide, img_file, title, desc, left, top,
                        img_w=Inches(3.8), suffix="", font_name="Calibri"):
    """Screenshot with title and description below."""
    add_screenshot(slide, img_file, left, top, img_w, suffix)
    add_text_box(slide, left, top + Inches(2.45), img_w, Inches(0.35),
                 title, 14, WHITE, True, PP_ALIGN.CENTER, font_name)
    add_text_box(slide, left, top + Inches(2.75), img_w, Inches(0.7),
                 desc, 10, LIGHT_GRAY, False, PP_ALIGN.CENTER, font_name)


# ─── SLIDE CONTENT DEFINITIONS ─────────────────────────────────────────

CONTENT = {
    "en": {
        "font": "Calibri",
        "suffix": "",
        "title": "510k Bridge",
        "subtitle": "FDA 510(k) Pathway Management\nfor Medical Device Companies",
        "tagline": "Pilot Software LLC dba 510k Bridge · Oregon, USA",
        "slides": {
            "problem_title": "The Challenge",
            "problem_sub": "Why medical device companies need a dedicated 510(k) partner",
            "problem_cards": [
                ("🏛", "Complex Regulation", "The FDA 510(k) process involves 17+ review criteria, 100+ guidance documents, and strict timeline requirements."),
                ("🌏", "Cross-Border Complexity", "International companies must navigate US entity formation, FDA registration, US Agent requirements, and labeling compliance."),
                ("⏱", "Timeline Pressure", "Average 510(k) review takes 130+ days. Delays in Pre-Sub, RTA screening, or AI/SI questions can add months."),
            ],
            "how_title": "How It Works",
            "how_sub": "Three steps from first call to FDA clearance",
            "how_steps": [
                ("1", "Schedule a Call", "Free 30-minute consultation to assess your device, predicate strategy, and timeline."),
                ("2", "Launch Control Tower", "We set up your tri-lingual dashboard, map milestones, and kick off the dual-track plan."),
                ("3", "Submit & Clear 510(k)", "We manage every gate, document, and FDA interaction through to clearance."),
            ],
            "services_title": "Services & Pricing",
            "services": [
                ("📊", "Control Tower SaaS", "From $500/mo", "Dedicated dashboard, tri-lingual wizard, dual-track tracking, document control, real-time messaging, risk & budget monitoring."),
                ("🏗️", "QMS-Lite", "$200–500/mo", "Lightweight QMS aligned to 21 CFR 820 & ISO 13485. Document control, CAPA, training, supplier qualification."),
                ("🔍", "Predicate Finder", "Included", "AI-powered openFDA search for predicate devices, substantial equivalence analysis, comparison reports."),
                ("⭐", "Professional PM", "$10–25K/mo", "Dedicated PMP manager, FDA Communications Center, US Agent, gate reviews, regulatory submission oversight."),
                ("🏛", "Entity Setup Tracker", "$1K–5K", "Delaware C-Corp, Oregon registration, EIN, bank account, FDA establishment, US Agent compliance."),
                ("🚀", "Enterprise", "Project-based", "Full service: regulatory strategy, 17-item RTA self-check, SE flowchart, US entity formation, end-to-end 510(k)."),
            ],
            "ct_title": "Inside the Control Tower",
            "ct_sub": "16-tab dashboard managing every aspect of your FDA journey",
            "showcase_slides": [
                {
                    "title": "Core Tracking",
                    "items": [
                        ("dual-track.png", "Dual-Track Dashboard", "Technical and regulatory milestones tracked side by side."),
                        ("gate-system.png", "Gate System", "Phase-gate reviews with Go/No-Go decisions and stakeholder inputs."),
                        ("risk-dashboard.png", "Risk Dashboard", "ISO 14971 risk matrix with severity, probability, and color-coded levels."),
                    ]
                },
                {
                    "title": "Documents & Compliance",
                    "items": [
                        ("document-control.png", "Document Control", "ISO 13485-aligned lifecycle tracking with version history and audit trail."),
                        ("audit-trail.png", "Audit Trail", "21 CFR Part 11 compliant — every change timestamped with full detail."),
                        ("fda-comms.png", "FDA Communications", "Q-Sub cover letters, 17-item RTA self-check, MDUFA timeline, SE flowchart."),
                    ]
                },
                {
                    "title": "Operations & Finance",
                    "items": [
                        ("budget.png", "Budget Tracking", "Planned vs. actual spend by category with automatic variance calculation."),
                        ("cash-runway.png", "Cash / Runway", "Burn rate analysis, runway projections, and funding milestone tracking."),
                        ("timeline.png", "Timeline", "Gantt-style project timeline with milestones and critical path."),
                    ]
                },
                {
                    "title": "Team & Collaboration",
                    "items": [
                        ("message-board.png", "Message Board", "Purpose-driven messaging tied to decisions, actions, and status updates."),
                        ("actions.png", "Actions / DHF / CAPA", "Task board, Design History File tracker, and Corrective Action log."),
                        ("suppliers.png", "Suppliers", "Supplier qualification, evaluation scoring, and document tracking."),
                    ]
                },
            ],
            "tools_title": "Standalone Tools",
            "tools_sub": "Powerful research tools included with every engagement",
            "tools_items": [
                ("predicate-finder.png", "510(k) Predicate Finder", "Search FDA's openFDA database by product code, applicant, or K-number. AI-powered substantial equivalence analysis."),
                ("guidance-docs.png", "FDA Guidance Search", "Search 600+ CDRH guidance documents — filter by topic, status, year. Direct links to official PDFs."),
                ("document-control.png", "QMS-Lite", "Lightweight QMS aligned to 21 CFR 820 and ISO 13485 — document control, CAPA, training, supplier qualification."),
            ],
            "highlights_title": "Platform Highlights",
            "highlights": [
                ("🧙", "Tri-lingual Wizard", "9-step wizard pre-populates your entire project in EN, 中文, or 한국어."),
                ("🏛️", "FDA Comms Center", "Q-Sub automation, RTA self-check, MDUFA tracking, SE flowchart — one tab."),
                ("📋", "RTA & DHF Readiness", "Real-time readiness scoring from your DHF and Standards trackers."),
                ("🌐", "Entity Setup", "Step-by-step US market entry — Delaware C-Corp, EIN, bank, FDA registration."),
            ],
            "resources_title": "Free Resources",
            "resources": [
                ("📘", "FDA 510(k) Pathway Guide", "Step-by-step US market entry guide"),
                ("📋", "Service Overview Fact Sheet", "One-page platform & capabilities summary"),
                ("🎓", "PMP Course", "Free project management course"),
                ("💰", "Investment Guide", "Term sheets, funding rounds, 510(k) Bridge strategy"),
            ],
            "contact_title": "Let's Talk",
            "contact_sub": "Schedule a free 30-minute consultation about your FDA pathway.",
            "contact_email": "info@510kbridge.com",
            "contact_web": "510kbridge.com",
            "contact_location": "Oregon, USA — Serving clients in US, China, and Korea",
            "contact_cta": "Schedule a Consultation →",
        }
    },
    "cn": {
        "font": "Calibri",
        "suffix": "-cn",
        "title": "510k Bridge",
        "subtitle": "中国医疗器械企业\nFDA 510(k) 通道管理",
        "tagline": "Pilot Software LLC（DBA 510k Bridge）· 俄勒冈州 & 上海",
        "slides": {
            "problem_title": "行业挑战",
            "problem_sub": "为什么医疗器械企业需要专业的510(k)合作伙伴",
            "problem_cards": [
                ("🏛", "法规复杂", "FDA 510(k)流程涉及17+审查标准、100+指导文件和严格的时间要求。"),
                ("🌏", "跨境复杂性", "国际企业需要完成美国实体注册、FDA注册、美国代理人要求和标签合规。"),
                ("⏱", "时间压力", "510(k)平均审查时间130+天。Pre-Sub、RTA筛查或AI/SI问题的延迟可能增加数月。"),
            ],
            "how_title": "如何运作",
            "how_sub": "从初次通话到FDA许可的三个步骤",
            "how_steps": [
                ("1", "预约咨询", "免费30分钟咨询，评估您的器械、等效器械策略和时间表。"),
                ("2", "启动 Control Tower", "我们为您搭建双语仪表盘，规划里程碑，启动双轨计划。"),
                ("3", "提交并获得510(k)许可", "我们管理每个关卡、文件和FDA互动，直至获得许可。"),
            ],
            "services_title": "服务与定价",
            "services": [
                ("📊", "Control Tower SaaS", "每月$500起", "专属仪表盘、双语向导、双轨跟踪、文档管控、实时消息、风险与预算监控。"),
                ("🏗️", "QMS-Lite 轻量版", "$200–500/月", "符合21 CFR 820和ISO 13485的轻量级QMS。文档控制、CAPA、培训、供应商资质。"),
                ("🔍", "等效器械查找器", "服务包含", "AI驱动的openFDA搜索，查找等效器械，实质性等效分析，对比报告。"),
                ("⭐", "专业项目管理", "$10,000–25,000/月", "专属PMP项目经理、FDA通信中心、美国代理人、关门评审、法规提交监督。"),
                ("🏛", "跨境实体设立追踪器", "$1K–5K", "特拉华C-Corp、俄勒冈登记、EIN、银行账户、FDA机构注册、美国代理人合规。"),
                ("🚀", "企业级", "项目定价", "全套服务：法规策略、17项RTA自检、SE流程图、美国实体成立、510(k)全流程管理。"),
            ],
            "ct_title": "深入了解 Control Tower",
            "ct_sub": "16个标签仪表盘管理FDA之旅的每一个方面",
            "showcase_slides": [
                {
                    "title": "核心跟踪",
                    "items": [
                        ("dual-track.png", "双轨道仪表盘", "技术和法规里程碑并排跟踪。"),
                        ("gate-system.png", "门控系统", "带标准清单的阶段门控评审和通过/不通过决定。"),
                        ("risk-dashboard.png", "风险仪表盘", "ISO 14971风险矩阵，含严重性、概率和颜色编码。"),
                    ]
                },
                {
                    "title": "文档与合规",
                    "items": [
                        ("document-control.png", "文档控制", "符合ISO 13485的生命周期跟踪，含版本历史和审计追溯。"),
                        ("audit-trail.png", "审计追踪", "符合21 CFR Part 11——每次更改都有时间戳和详细记录。"),
                        ("fda-comms.png", "FDA通讯中心", "Q-Sub附信、17项RTA自检、MDUFA时间线、SE流程图。"),
                    ]
                },
                {
                    "title": "运营与财务",
                    "items": [
                        ("budget.png", "预算跟踪", "按类别的计划与实际支出及自动差异计算。"),
                        ("cash-runway.png", "资金 / 跑道", "燃烧率分析、跑道预测和融资里程碑跟踪。"),
                        ("timeline.png", "时间线", "甘特图风格的项目时间线，含里程碑和关键路径。"),
                    ]
                },
                {
                    "title": "团队与协作",
                    "items": [
                        ("message-board.png", "消息中心", "目标驱动的消息系统，与决策和行动紧密结合。"),
                        ("actions.png", "行动项 / DHF / CAPA", "任务板、设计历史文件和纠正预防措施日志。"),
                        ("suppliers.png", "供应商", "供应商资质、评估评分和文档跟踪。"),
                    ]
                },
            ],
            "tools_title": "独立工具",
            "tools_sub": "每次合作都包含强大的研究工具",
            "tools_items": [
                ("predicate-finder.png", "510(k) 等效器械查找器", "搜索FDA openFDA数据库，AI驱动的实质性等效分析。"),
                ("guidance-docs.png", "FDA指导文件搜索", "搜索600+ CDRH指导文件——按主题、状态、年份筛选。"),
                ("document-control.png", "QMS-Lite 轻量版", "符合21 CFR 820和ISO 13485的轻量级质量管理系统。"),
            ],
            "highlights_title": "平台亮点",
            "highlights": [
                ("🧙", "三语设置向导", "9步向导预填充整个项目——英文、中文或韩文。"),
                ("🏛️", "FDA通信中心", "Q-Sub自动化、RTA自检、MDUFA跟踪、SE流程图。"),
                ("📋", "RTA与DHF就绪评估", "从DHF和标准跟踪器实时评估就绪状态。"),
                ("🌐", "跨境实体设立", "美国市场准入——特拉华C-Corp、EIN、银行、FDA注册。"),
            ],
            "resources_title": "免费资源",
            "resources": [
                ("📘", "FDA 510(k)通道指南", "逐步美国市场准入指南"),
                ("📋", "服务概览单页", "平台和功能一页概览"),
                ("🎓", "PMP课程", "免费项目管理课程"),
                ("💰", "投资融资指南", "条款清单、融资轮次、510(k) Bridge策略"),
            ],
            "contact_title": "联系我们",
            "contact_sub": "预约免费30分钟咨询，讨论您的FDA通道。",
            "contact_email": "info@510kbridge.com",
            "contact_web": "510kbridge.com",
            "contact_location": "俄勒冈州，美国 — 服务中国、韩国和美国客户",
            "contact_cta": "预约咨询 →",
        }
    },
}


def build_presentation(lang):
    c = CONTENT[lang]
    s = c["slides"]
    fn = c["font"]
    sx = c["suffix"]
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank

    # ── SLIDE 1: Title ──────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 c["title"], 56, ACCENT, True, PP_ALIGN.CENTER, fn)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.2),
                 c["subtitle"], 28, WHITE, False, PP_ALIGN.CENTER, fn)
    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
                 c["tagline"], 16, LIGHT_GRAY, False, PP_ALIGN.CENTER, fn)

    # ── SLIDE 2: The Problem ────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["problem_title"], font_name=fn)
    add_subtitle_text(slide, s["problem_sub"], font_name=fn)
    for i, (icon, title, body) in enumerate(s["problem_cards"]):
        add_card(slide, Inches(0.8 + i * 4.1), Inches(2.0),
                 Inches(3.8), Inches(4.0), icon, title, body, fn)

    # ── SLIDE 3: How It Works ───────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["how_title"], font_name=fn)
    add_subtitle_text(slide, s["how_sub"], font_name=fn)
    for i, (num, title, body) in enumerate(s["how_steps"]):
        x = Inches(0.8 + i * 4.1)
        # Step number circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.4),
                                       Inches(2.2), Inches(1.0), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK_BG
        p.font.name = fn
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_text_box(slide, x, Inches(3.5), Inches(3.8), Inches(0.5),
                     title, 20, WHITE, True, PP_ALIGN.CENTER, fn)
        add_text_box(slide, x, Inches(4.1), Inches(3.8), Inches(1.5),
                     body, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER, fn)
        # Arrow between steps
        if i < 2:
            add_text_box(slide, x + Inches(3.9), Inches(2.4), Inches(0.4),
                         Inches(0.6), "→", 32, ACCENT, True, PP_ALIGN.CENTER, fn)

    # ── SLIDE 4: Services ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["services_title"], font_name=fn)
    for i, (icon, title, price, body) in enumerate(s["services"]):
        row = i // 3
        col = i % 3
        x = Inches(0.6 + col * 4.1)
        y = Inches(1.4 + row * 3.0)
        card_text = f"{title}\n{price}"
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       x, y, Inches(3.8), Inches(2.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(10)
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)

        p0 = tf.paragraphs[0]
        p0.text = f"{icon}  {title}"
        p0.font.size = Pt(16)
        p0.font.bold = True
        p0.font.color.rgb = WHITE
        p0.font.name = fn

        p1 = tf.add_paragraph()
        p1.text = price
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = GOLD
        p1.font.name = fn
        p1.space_before = Pt(4)

        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(10)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = fn
        p2.space_before = Pt(6)

    # ── SLIDE 5: Control Tower Intro ────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["ct_title"], font_name=fn)
    add_subtitle_text(slide, s["ct_sub"], font_name=fn)
    # Large hero screenshot
    add_screenshot(slide, "dual-track.png", Inches(1.5), Inches(2.0),
                   Inches(10.3), sx)

    # ── SLIDES 6-9: Showcase groups ─────────────────────────────────
    for group in s["showcase_slides"]:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide)
        add_title_text(slide, f"{s['ct_title']} — {group['title']}", 32, font_name=fn)
        for i, (img, title, desc) in enumerate(group["items"]):
            add_screenshot_card(slide, img, title, desc,
                                Inches(0.5 + i * 4.2), Inches(1.6),
                                Inches(3.9), sx, fn)

    # ── SLIDE 10: Standalone Tools ──────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["tools_title"], font_name=fn)
    add_subtitle_text(slide, s["tools_sub"], font_name=fn)
    for i, (img, title, desc) in enumerate(s["tools_items"]):
        add_screenshot_card(slide, img, title, desc,
                            Inches(0.5 + i * 4.2), Inches(2.0),
                            Inches(3.9), sx, fn)

    # ── SLIDE 11: Platform Highlights ───────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["highlights_title"], font_name=fn)
    for i, (icon, title, body) in enumerate(s["highlights"]):
        col = i % 2
        row = i // 2
        add_card(slide, Inches(0.8 + col * 6.2), Inches(1.6 + row * 2.8),
                 Inches(5.8), Inches(2.4), icon, title, body, fn)

    # ── SLIDE 12: Free Resources ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_title_text(slide, s["resources_title"], font_name=fn)
    for i, (icon, title, desc) in enumerate(s["resources"]):
        add_card(slide, Inches(0.4 + i * 3.2), Inches(1.8),
                 Inches(2.9), Inches(3.5), icon, title, desc, fn)

    # ── SLIDE 13: Contact ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.0),
                 s["contact_title"], 44, ACCENT, True, PP_ALIGN.CENTER, fn)
    add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(0.6),
                 s["contact_sub"], 20, LIGHT_GRAY, False, PP_ALIGN.CENTER, fn)

    # Contact details card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(3.5), Inches(3.8), Inches(6.3), Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(20)
    tf.margin_left = Pt(24)

    lines = [
        (f"📧  {s['contact_email']}", 18, WHITE),
        (f"🌐  {s['contact_web']}", 18, WHITE),
        (f"📍  {s['contact_location']}", 16, LIGHT_GRAY),
        ("", 12, WHITE),
        (s["contact_cta"], 22, ACCENT),
    ]
    for j, (text, size, color) in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = (j == len(lines) - 1)
        p.font.name = fn
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

    return prs


if __name__ == "__main__":
    os.makedirs(OUT_DIR, exist_ok=True)

    for lang, label in [("en", "EN"), ("cn", "CN")]:
        prs = build_presentation(lang)
        out = os.path.join(OUT_DIR, f"510kBridge_Presentation_{label}.pptx")
        prs.save(out)
        print(f"✅ {out}")
